import re
import fitz  # PyMuPDF for PDF processing
from pdf2image import convert_from_bytes
from PIL import Image
import io
from io import BytesIO
import easyocr
import numpy as np
import ssl
from pptx import Presentation
from docx import Document
from multiprocessing import Pool, cpu_count

# Create an unverified SSL context
ssl._create_default_https_context = ssl._create_unverified_context

# Initialize EasyOCR reader (English as the language)
reader = easyocr.Reader(['en'], gpu=True)


def extract_text_and_images(file):
    file_type = file.type  # Updated: Get type from the UploadedFile object
    try:
        file_content = file.read()
        if file_type == "application/pdf":
            return extract_text_from_pdf(file_content)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return extract_text_from_ppt(file_content)
        elif file_type == "text/plain":
            return extract_text_from_text(file_content)
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return extract_text_from_doc(file_content)
    except Exception as e:
        print(f"Error processing file: {e}")
        raise RuntimeError(f"Failed to process file: {e}")


def extract_text_from_pdf(file):
    """Extracts text and images from an uploaded PDF using EasyOCR."""
    text_content = []

    # Open the PDF from the uploaded file-like object
    with fitz.open(stream=file, filetype="pdf") as pdf:
        for page_num in range(pdf.page_count):
            page = pdf.load_page(page_num)

            # Extract text from the PDF page
            text = page.get_text()
            text_content.append(text)

            # Extract images from the PDF page
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = pdf.extract_image(xref)
                image_bytes = base_image.get("image")

                # Skip if image data is missing or corrupted
                if not image_bytes:
                    continue

                # Convert image bytes to a PIL Image and validate
                try:
                    image = Image.open(io.BytesIO(image_bytes))
                    image_np = np.array(image)  # Convert to numpy array

                    # Ensure the image is valid (non-empty)
                    if image_np.size == 0:
                        continue  # Ignoring empty images , no need for raising error

                    # Skip if the image dimensions are invalid
                    if image_np.size == 0 or image_np.shape[0] == 0 or image_np.shape[1] == 0:
                        continue

                    # Apply OCR on the image using EasyOCR
                    ocr_results = reader.readtext(image_np)

                    # Collect OCR results
                    extracted_text = "\n".join([res[1] for res in ocr_results])
                    text_content.append(extracted_text)

                except Exception:
                    # Skip this image if any error occurs during processing
                    continue

    text_content = clean_extracted_text(text_content)
    return text_content


def extract_text_from_ppt(file_content):
    """Extracts text from an uploaded PPTX file."""
    text_content = []
    try:
        # Wrap bytes in BytesIO
        ppt_stream = BytesIO(file_content)
        prs = Presentation(ppt_stream)

        for slide_num, slide in enumerate(prs.slides):
            slide_content = []  # Temporary list for content from this slide
            try:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_content.append(run.text)

                    # Process images in the slide
                    if shape.shape_type == 13:  # Check if the shape is an image
                        try:
                            image_bytes = shape.image_blob
                            image = Image.open(io.BytesIO(image_bytes))
                            image_np_array = np.array(image)

                            if image_np_array.size > 0:
                                ocr_response = reader.readtext(image_np_array)
                                for img in ocr_response:
                                    slide_content.append(img[1])
                        except Exception as e:
                            print(
                                f"Error processing image on slide {slide_num + 1}: {e}"
                            )
            except Exception as e:
                print(
                    f"Error occurred while processing slide {slide_num + 1}: {e}")

            # Add slide content to the final list
            text_content.append("\n".join(slide_content))

    except Exception as e:
        raise RuntimeError(f"Error processing PPTX file: {e}")

    text_content = clean_extracted_text(text_content)
    return text_content


def extract_text_from_doc(file):
    """Extracts text form an uploaded document file. """
    text_content = []
    try:
        doc = Document(file)
        for para in doc.paragraphs:
            text_content.append(para.text)

        for table in doc.tables:
            for row in table.rows:
                for col in row.cells:
                    text_content.append(col.text)

        for i in doc.part.rels.values():
            if i.reltype == 13:
                image_data = doc.part.get_blob(i.target)
                image = Image.open(io.BytesIO(image_data))
                image_np_array = np.array(image)

                if image_np_array.size == 0:
                    continue  # Ignoring empty images , no need for raising error

                ocr_response = reader.readtext(image_np_array)
                text_extracted = "\n".join([i[1] for i in ocr_response])
                text_content.append(text_extracted)
    except Exception as e:
        raise RuntimeError(f"Error processing text file: {e}")

    text_content = clean_extracted_text(text_content)
    return text_content


def extract_text_from_text(file):
    """Extracts text from an uploaded text file."""
    text_content = []
    try:
        text_content = file.decode('utf-8').splitlines()
    except UnicodeDecodeError:
        file.seek(0)
        text_content = file.decode('latin-1').splitlines()
    except Exception as e:
        raise RuntimeError(f"Error processing document file: {e}")
    text_content = clean_extracted_text(text_content)
    return text_content


def clean_extracted_text(text):
    """
    Cleans the extracted text by removing irrelevant words/phrases (stopwords).

    Args:
        text (list): List of text lines extracted from the document.
        stopwords (list): List of words/phrases to exclude.

    Returns:
        list: Cleaned text content with stopwords removed.
    """
    stopwords = ["Purdue", "Outline", "Chapter", "University", "Fort Wayne", "CONTD", "DEMO", "Q & A", "Have Fun",
                 "PURDUE", "purdue", "UNIVERSITY", "U N [ V E R $ [ T Y", "FORT WAYNE", "U N I V E R $ I T Y", "U N I V E R S I T Y"]

    cleaned_text = []
    for line in text:
        # Remove each stopword from the line
        for stopword in stopwords:
            line = line.replace(stopword, "").strip()

        # Remove extra whitespace and skip empty lines
        if line:
            cleaned_text.append(line)

    return cleaned_text
